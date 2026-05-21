"""
16-slide PPTX builder using SENSUM brand colors and fonts.

All charts are rendered via matplotlib to temp PNGs, then embedded.
"""

import tempfile
from collections import Counter
from pathlib import Path

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------- #
# Brand constants
# --------------------------------------------------------------------------- #

BEIGE = RGBColor(0xF4, 0xE5, 0xCA)
BROWN = RGBColor(0x58, 0x2F, 0x0E)
BEIGE_HEX = "#F4E5CA"
BROWN_HEX = "#582F0E"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

FONT_TITLE = "EB Garamond"
FONT_BODY = "Lora"

PROJECT_ROOT = Path(__file__).parent.parent.parent
FONT_DIR = PROJECT_ROOT / "outputs" / "channel_assets" / "fonts"
FONT_TITLE_PATH = FONT_DIR / "EBGaramond-Bold.ttf"
FONT_BODY_PATH = FONT_DIR / "Lora-Italic.ttf"

DAYS = ["Mon", "Tue", "Wed", "Thu", "Fri", "Sat", "Sun"]


def _register_fonts():
    for p in (FONT_TITLE_PATH, FONT_BODY_PATH):
        if p.exists():
            fm.fontManager.addfont(str(p))


_register_fonts()


def _title_fp():
    if FONT_TITLE_PATH.exists():
        return fm.FontProperties(fname=str(FONT_TITLE_PATH))
    return fm.FontProperties(family="serif")


def _body_fp():
    if FONT_BODY_PATH.exists():
        return fm.FontProperties(fname=str(FONT_BODY_PATH))
    return fm.FontProperties(family="serif", style="italic")


# --------------------------------------------------------------------------- #
# Presentation helpers
# --------------------------------------------------------------------------- #

def _new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def _blank_slide(prs: Presentation):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BEIGE
    return slide


def _add_text(slide, text: str, left, top, width, height,
              font_name=FONT_TITLE, size=28, bold=False, align=PP_ALIGN.LEFT,
              color=None):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or BROWN
    return txb


def _add_rule(slide, top, left=Inches(0.5), width=Inches(12.333), thickness=Pt(0.75)):
    line = slide.shapes.add_connector(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT
        left, top,
        left + width, top,
    )
    line.line.color.rgb = BROWN
    line.line.width = thickness


def _slide_header(slide, title: str, slide_num: int):
    _add_text(slide, title,
              left=Inches(0.5), top=Inches(0.3),
              width=Inches(11.5), height=Inches(0.8),
              font_name=FONT_TITLE, size=30, bold=True)
    _add_text(slide, str(slide_num),
              left=Inches(12.3), top=Inches(0.3),
              width=Inches(0.8), height=Inches(0.5),
              font_name=FONT_BODY, size=11, align=PP_ALIGN.RIGHT,
              color=RGBColor(0x99, 0x7A, 0x5C))
    _add_rule(slide, top=Inches(1.15))


def _embed_chart(slide, fig, left, top, width, height):
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
        fig.savefig(tmp.name, dpi=150, bbox_inches="tight",
                    facecolor=BEIGE_HEX, edgecolor="none")
        tmp_path = tmp.name
    plt.close(fig)
    slide.shapes.add_picture(tmp_path, left, top, width, height)
    Path(tmp_path).unlink(missing_ok=True)


def _mpl_fig(w_in=10.5, h_in=5.5):
    fig, ax = plt.subplots(figsize=(w_in, h_in))
    fig.patch.set_facecolor(BEIGE_HEX)
    ax.set_facecolor(BEIGE_HEX)
    for spine in ax.spines.values():
        spine.set_color(BROWN_HEX)
        spine.set_linewidth(0.8)
    ax.tick_params(colors=BROWN_HEX, labelsize=9)
    ax.xaxis.label.set_color(BROWN_HEX)
    ax.yaxis.label.set_color(BROWN_HEX)
    return fig, ax


def _apply_title_fp(ax, title="", xlabel="", ylabel=""):
    fp = _title_fp()
    if title:
        ax.set_title(title, fontproperties=fp, color=BROWN_HEX, fontsize=11, pad=8)
    if xlabel:
        ax.set_xlabel(xlabel, fontproperties=_body_fp(), color=BROWN_HEX, fontsize=9)
    if ylabel:
        ax.set_ylabel(ylabel, fontproperties=_body_fp(), color=BROWN_HEX, fontsize=9)


def _tick_font(ax):
    fp = _body_fp()
    for label in ax.get_xticklabels() + ax.get_yticklabels():
        label.set_fontproperties(fp)
        label.set_color(BROWN_HEX)
        label.set_fontsize(8)


# --------------------------------------------------------------------------- #
# Individual slide builders
# --------------------------------------------------------------------------- #

def _slide_01_cover(prs, week_label: str):
    slide = _blank_slide(prs)
    _add_text(slide, "SENSUM",
              Inches(0.5), Inches(1.8), Inches(12), Inches(1.2),
              font_name=FONT_TITLE, size=60, bold=True, align=PP_ALIGN.CENTER)
    _add_text(slide, "Niche Intelligence Report",
              Inches(0.5), Inches(3.0), Inches(12), Inches(0.8),
              font_name=FONT_BODY, size=24, align=PP_ALIGN.CENTER,
              color=RGBColor(0x7A, 0x4A, 0x20))
    _add_rule(slide, top=Inches(3.9), left=Inches(3.5), width=Inches(6))
    _add_text(slide, week_label,
              Inches(0.5), Inches(4.1), Inches(12), Inches(0.6),
              font_name=FONT_BODY, size=16, align=PP_ALIGN.CENTER,
              color=RGBColor(0x7A, 0x4A, 0x20))
    _add_text(slide, "Psychology & Mental Health · Medium Scope",
              Inches(0.5), Inches(4.7), Inches(12), Inches(0.5),
              font_name=FONT_BODY, size=11, align=PP_ALIGN.CENTER,
              color=RGBColor(0xAA, 0x80, 0x55))


def _slide_02_key_findings(prs, results: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Key Findings", 2)

    findings = _derive_key_findings(results)
    top = Inches(1.4)
    for i, f in enumerate(findings[:5]):
        _add_text(slide, f"— {f}",
                  Inches(0.7), top + Inches(i * 1.0),
                  Inches(12), Inches(0.9),
                  font_name=FONT_BODY, size=14)


def _derive_key_findings(results: dict) -> list[str]:
    findings = []
    channels = results.get("channels", {})
    videos = results.get("videos", [])
    topics = results.get("trending_topics", Counter())
    gaps = results.get("gaps", [])

    top_ch = sorted(channels.values(), key=lambda c: c.get("subscribers", 0), reverse=True)
    if top_ch:
        findings.append(
            f"Top channel this week: {top_ch[0]['name']} "
            f"({top_ch[0]['subscribers']:,} subscribers)"
        )

    velocities = results.get("view_velocity", [])
    if velocities:
        v = velocities[0]
        findings.append(
            f"Fastest-rising video: \"{v['title'][:60]}\" "
            f"({int(v['velocity']):,} views/day)"
        )

    if topics:
        top3 = ", ".join(list(topics.keys())[:3])
        findings.append(f"Top trending keywords: {top3}")

    if gaps:
        findings.append(
            f"Content gap opportunity: \"{gaps[0]['topic']}\" "
            f"(high niche engagement, absent from SENSUM)"
        )

    non_sensum = [v for v in videos if not channels.get(v.get("channel_id", {}), {}).get("is_sensum")]
    if non_sensum:
        eng = results.get("engagement", {})
        top_eng = sorted(non_sensum, key=lambda v: eng.get(v["video_id"], 0), reverse=True)
        if top_eng:
            findings.append(
                f"Highest engagement: \"{top_eng[0]['title'][:55]}\" "
                f"({eng.get(top_eng[0]['video_id'], 0):.2f}% rate)"
            )
    return findings


def _slide_03_channel_subscribers(prs, channels: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Top Channels by Subscribers", 3)

    sorted_ch = sorted(channels.values(), key=lambda c: c.get("subscribers", 0), reverse=True)[:12]
    names = [c["name"][:25] for c in sorted_ch]
    subs = [c.get("subscribers", 0) for c in sorted_ch]

    fig, ax = _mpl_fig(10.5, 5.5)
    bars = ax.barh(names[::-1], subs[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Subscribers", fontproperties=_body_fp())
    for bar, val in zip(bars, subs[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                f"{val:,}", va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda x, _: f"{int(x/1000)}k" if x >= 1000 else str(int(x))))
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_04_channel_growth(prs, channels: dict, prev_snapshot: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Channel Growth — Week over Week", 4)

    deltas = {}
    for cid, ch in channels.items():
        prev = prev_snapshot.get(cid)
        if prev is not None:
            deltas[ch["name"][:25]] = ch.get("subscribers", 0) - prev
        else:
            deltas[ch["name"][:25]] = 0

    sorted_d = sorted(deltas.items(), key=lambda x: abs(x[1]), reverse=True)[:10]
    names = [n for n, _ in sorted_d]
    vals = [v for _, v in sorted_d]

    fig, ax = _mpl_fig(10.5, 5.5)
    colors = [BROWN_HEX if v >= 0 else "#C0956A" for v in vals]
    ax.barh(names[::-1], vals[::-1], color=colors[::-1], alpha=0.85)
    ax.axvline(0, color=BROWN_HEX, linewidth=0.8, alpha=0.5)
    ax.set_xlabel("Subscriber delta", fontproperties=_body_fp())
    if not any(vals):
        ax.text(0.5, 0.5, "No prior week data yet.\nGrowth trends appear from week 2 onward.",
                transform=ax.transAxes, ha="center", va="center",
                fontproperties=_body_fp(), fontsize=12, color=BROWN_HEX, alpha=0.6)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_05_channel_table(prs, channels: dict, engagement: dict, videos: list[dict]):
    slide = _blank_slide(prs)
    _slide_header(slide, "Channel Profiles", 5)

    ch_eng: dict[str, list[float]] = {}
    for v in videos:
        cid = v["channel_id"]
        rate = engagement.get(v["video_id"], 0.0)
        ch_eng.setdefault(cid, []).append(rate)

    sorted_ch = sorted(channels.values(), key=lambda c: c.get("subscribers", 0), reverse=True)[:10]

    rows = []
    for ch in sorted_ch:
        cid = ch["channel_id"]
        avg_eng = sum(ch_eng.get(cid, [0])) / max(len(ch_eng.get(cid, [1])), 1)
        sensum_mark = " ★" if ch.get("is_sensum") else ""
        rows.append([
            ch["name"][:22] + sensum_mark,
            f"{ch.get('subscribers', 0):,}",
            f"{ch.get('video_count', 0):,}",
            f"{avg_eng:.2f}%",
        ])

    headers = ["Channel", "Subscribers", "Total Videos", "Avg Engagement"]
    col_widths = [Inches(3.8), Inches(2.2), Inches(2.2), Inches(2.5)]
    col_x = [Inches(0.5), Inches(4.3), Inches(6.5), Inches(8.7)]
    row_h = Inches(0.47)
    top_start = Inches(1.35)

    # Header row
    for j, (h, x, w) in enumerate(zip(headers, col_x, col_widths)):
        _add_text(slide, h, x, top_start, w, row_h,
                  font_name=FONT_TITLE, size=11, bold=True)
    _add_rule(slide, top=top_start + row_h - Inches(0.05),
              left=Inches(0.5), width=Inches(11.0), thickness=Pt(1.2))

    # Data rows
    for i, row in enumerate(rows):
        y = top_start + row_h + Inches(i * 0.47)
        bg_alpha = 0.04 if i % 2 == 0 else 0
        for j, (cell, x, w) in enumerate(zip(row, col_x, col_widths)):
            _add_text(slide, cell, x, y, w, row_h,
                      font_name=FONT_BODY, size=10)

    _add_text(slide, "★ = SENSUM",
              Inches(0.5), Inches(7.1), Inches(3), Inches(0.3),
              font_name=FONT_BODY, size=8,
              color=RGBColor(0x99, 0x7A, 0x5C))


def _slide_06_trending_topics(prs, topics: Counter):
    slide = _blank_slide(prs)
    _slide_header(slide, "Trending Topics", 6)

    if not topics:
        _add_text(slide, "No data yet.", Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    items = topics.most_common(20)
    labels, vals = zip(*items)

    fig, ax = _mpl_fig(10.5, 5.5)
    ax.barh(list(labels)[::-1], list(vals)[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Frequency in video titles", fontproperties=_body_fp())
    for bar, val in zip(ax.patches, list(vals)[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _fallback_bar(slide, counter: Counter, xlabel="", ylabel=""):
    items = counter.most_common(15)
    labels, vals = zip(*items) if items else ([], [])
    fig, ax = _mpl_fig(10.5, 5.2)
    ax.barh(list(labels)[::-1], list(vals)[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel(ylabel, fontproperties=_body_fp())
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_07_top_videos(prs, velocity_list: list[dict]):
    slide = _blank_slide(prs)
    _slide_header(slide, "Top Videos This Week (View Velocity)", 7)

    top = velocity_list[:10]
    if not top:
        return

    labels = [f"{v['title'][:48]}…" if len(v['title']) > 48 else v['title'] for v in top]
    vals = [v["velocity"] for v in top]

    fig, ax = _mpl_fig(12, 5.5)
    ax.barh(labels[::-1], vals[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Views per day since publish", fontproperties=_body_fp())
    for bar, val in zip(ax.patches, vals[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                f"{val:,.0f}", va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    fig.subplots_adjust(left=0.42, right=0.88, top=0.97, bottom=0.12)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_08_channel_engagement(prs, videos: list[dict], channels: dict, engagement: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Top Channels by Engagement Rate", 8)

    ch_eng: dict[str, list[float]] = {}
    for v in videos:
        cid = v["channel_id"]
        rate = engagement.get(v["video_id"], 0.0)
        if rate > 0:
            ch_eng.setdefault(cid, []).append(rate)

    ch_avgs = {}
    for cid, rates in ch_eng.items():
        if len(rates) >= 2:
            ch_name = channels.get(cid, {}).get("name", cid)[:25]
            ch_avgs[ch_name] = round(sum(rates) / len(rates), 3)

    if not ch_avgs:
        _add_text(slide, "Not enough video data for engagement analysis.",
                  Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    sorted_ch = sorted(ch_avgs.items(), key=lambda x: x[1], reverse=True)[:12]
    names = [n for n, _ in sorted_ch]
    vals = [v for _, v in sorted_ch]

    fig, ax = _mpl_fig(10.5, 5.5)
    bars = ax.barh(names[::-1], vals[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Avg engagement rate (%)", fontproperties=_body_fp())
    for bar, val in zip(bars, vals[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                f"{val:.2f}%", va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_09_duration(prs, dur_eng: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Optimal Video Length", 9)

    labels = list(dur_eng.keys())
    vals = list(dur_eng.values())

    fig, ax = _mpl_fig(8, 5)
    bars = ax.bar(labels, vals, color=BROWN_HEX, alpha=0.85, width=0.55)
    ax.set_ylabel("Avg engagement rate (%)", fontproperties=_body_fp())
    for bar, val in zip(bars, vals):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 0.002,
                f"{val:.2f}%", ha="center", fontproperties=_body_fp(),
                fontsize=9, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(1.5), Inches(1.3), Inches(10.3), Inches(6.0))


def _slide_10_content_gaps(prs, gaps: list[dict]):
    slide = _blank_slide(prs)
    _slide_header(slide, "Content Gaps", 10)

    _add_text(slide,
              "Topics with high niche engagement absent from SENSUM's published corpus:",
              Inches(0.5), Inches(1.25), Inches(12), Inches(0.5),
              font_name=FONT_BODY, size=11, color=RGBColor(0x7A, 0x4A, 0x20))

    for i, gap in enumerate(gaps[:8]):
        score_bar = "█" * min(int(gap["engagement_score"] / max((gaps[0]["engagement_score"] or 1), 1) * 12), 12)
        line = f"{gap['topic']}   {score_bar}  ({gap['engagement_score']:,})"
        _add_text(slide, line,
                  Inches(0.7), Inches(1.85) + Inches(i * 0.63),
                  Inches(11.5), Inches(0.55),
                  font_name=FONT_BODY, size=13)


def _slide_11_recommendations(prs, gaps: list[dict], topics: Counter):
    slide = _blank_slide(prs)
    _slide_header(slide, "This Week's Topic Recommendations", 11)

    _add_text(slide,
              "Three SENSUM-angle ideas derived from gap + trending data:",
              Inches(0.5), Inches(1.25), Inches(12), Inches(0.5),
              font_name=FONT_BODY, size=11, color=RGBColor(0x7A, 0x4A, 0x20))

    top_gaps = [g["topic"] for g in gaps[:6]]
    top_topics = [w for w in topics.keys() if w not in top_gaps][:6]

    ideas = []
    for g in top_gaps[:3]:
        ideas.append(f"The shame around {g} — what the research actually says")
    for t in top_topics[:max(0, 3 - len(ideas))]:
        ideas.append(f"Why {t} feels so wrong (and why it isn't)")

    positions = [Inches(2.0), Inches(3.7), Inches(5.4)]
    for i, idea in enumerate(ideas[:3]):
        y = positions[i]
        _add_text(slide, f"{i+1}.", Inches(0.5), y, Inches(0.4), Inches(0.7),
                  font_name=FONT_TITLE, size=28, bold=True,
                  color=RGBColor(0xAA, 0x80, 0x55))
        _add_text(slide, idea, Inches(1.1), y, Inches(11.5), Inches(0.9),
                  font_name=FONT_BODY, size=16)


def _slide_12_title_words(prs, title_words: Counter):
    slide = _blank_slide(prs)
    _slide_header(slide, "Emotional Trigger Words in Titles", 12)

    if not title_words:
        _add_text(slide, "No emotional trigger words found in current dataset.",
                  Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    items = title_words.most_common(20)
    labels, vals = zip(*items)

    fig, ax = _mpl_fig(10.5, 5.5)
    ax.barh(list(labels)[::-1], list(vals)[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Occurrences in video titles", fontproperties=_body_fp())
    for bar, val in zip(ax.patches, list(vals)[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_13_comment_sentiment(prs, sentiment: Counter):
    slide = _blank_slide(prs)
    _slide_header(slide, "Comment Sentiment Clusters", 13)

    if not sentiment:
        _add_text(slide, "No comment data available for this run.",
                  Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    items = sentiment.most_common(15)
    labels, vals = zip(*items)

    fig, ax = _mpl_fig(10.5, 5.5)
    alphas = [0.9 - 0.04 * i for i in range(len(labels))]
    for i, (label, val, alpha) in enumerate(zip(labels[::-1], vals[::-1], alphas[::-1])):
        ax.barh(label, val, color=BROWN_HEX, alpha=alpha)
    ax.set_xlabel("Occurrences in top-video comments", fontproperties=_body_fp())
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


def _slide_14_thumbnail_styles(prs, videos: list[dict]):
    slide = _blank_slide(prs)
    _slide_header(slide, "Thumbnail Style Breakdown", 14)

    counts: Counter = Counter()
    for v in videos:
        style = v.get("thumbnail_style", "unknown")
        if style and style != "unknown":
            counts[style] += 1

    if not counts:
        _add_text(slide, "Thumbnail classification not yet available.",
                  Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    labels = list(counts.keys())
    sizes = list(counts.values())
    alphas = [0.9, 0.7, 0.55, 0.4, 0.25]

    fig, ax = plt.subplots(figsize=(8, 5.5))
    fig.patch.set_facecolor(BEIGE_HEX)
    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, autopct="%1.0f%%",
        colors=[f"rgba({int(0x58)},{int(0x2F)},{int(0x0E)},{a})"
                if False  # matplotlib doesn't take rgba strings
                else BROWN_HEX for a in alphas[:len(sizes)]],
        startangle=90,
        wedgeprops=dict(linewidth=1.5, edgecolor=BEIGE_HEX),
        pctdistance=0.8,
    )
    # Re-color wedges with different opacities via face color
    from matplotlib.patches import Wedge
    for i, wedge in enumerate(wedges):
        alpha = alphas[i % len(alphas)]
        wedge.set_facecolor(BROWN_HEX)
        wedge.set_alpha(alpha)

    fp = _body_fp()
    for t in texts + autotexts:
        t.set_color(BROWN_HEX)
        t.set_fontproperties(fp)
        t.set_fontsize(10)

    ax.set_facecolor(BEIGE_HEX)
    plt.tight_layout(pad=0.3)
    _embed_chart(slide, fig, Inches(2.5), Inches(1.3), Inches(8.3), Inches(6.0))


def _slide_15_publish_timing(prs, timing_matrix: list[list[int]]):
    slide = _blank_slide(prs)
    _slide_header(slide, "Publish Timing Heatmap", 15)

    import numpy as np
    data = np.array(timing_matrix)  # 7 × 24

    fig, ax = _mpl_fig(11, 4.8)
    im = ax.imshow(data, aspect="auto", cmap="YlOrBr",
                   interpolation="nearest", vmin=0)
    ax.set_xticks(range(24))
    ax.set_xticklabels([f"{h:02d}h" if h % 3 == 0 else "" for h in range(24)],
                       fontsize=7, color=BROWN_HEX)
    ax.set_yticks(range(7))
    ax.set_yticklabels(DAYS, fontproperties=_body_fp(), fontsize=9, color=BROWN_HEX)
    ax.set_xlabel("Hour (UTC)", fontproperties=_body_fp(), fontsize=9, color=BROWN_HEX)

    cbar = fig.colorbar(im, ax=ax, shrink=0.8, pad=0.02)
    cbar.ax.tick_params(colors=BROWN_HEX, labelsize=7)
    cbar.set_label("Uploads", color=BROWN_HEX, fontsize=8)

    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(5.8))


def _slide_16_upload_frequency(prs, videos: list[dict], channels: dict):
    slide = _blank_slide(prs)
    _slide_header(slide, "Upload Frequency (Last 30 Days)", 16)

    from collections import Counter as _Counter
    counts = _Counter(v["channel_id"] for v in videos)

    named = {
        channels.get(cid, {}).get("name", cid)[:25]: cnt
        for cid, cnt in counts.items()
        if cid in channels
    }

    if not named:
        _add_text(slide, "No video data available.",
                  Inches(0.5), Inches(2), Inches(12), Inches(1),
                  font_name=FONT_BODY, size=14)
        return

    sorted_data = sorted(named.items(), key=lambda x: x[1], reverse=True)[:15]
    names = [n for n, _ in sorted_data]
    vals = [v for _, v in sorted_data]

    fig, ax = _mpl_fig(10.5, 5.5)
    bars = ax.barh(names[::-1], vals[::-1], color=BROWN_HEX, alpha=0.85)
    ax.set_xlabel("Videos published in last 30 days", fontproperties=_body_fp())
    for bar, val in zip(bars, vals[::-1]):
        ax.text(bar.get_width() * 1.01, bar.get_y() + bar.get_height() / 2,
                str(val), va="center", fontproperties=_body_fp(),
                fontsize=8, color=BROWN_HEX)
    _apply_title_fp(ax)
    _tick_font(ax)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)
    plt.tight_layout(pad=0.5)
    _embed_chart(slide, fig, Inches(0.4), Inches(1.3), Inches(12.5), Inches(6.0))


# --------------------------------------------------------------------------- #
# Main entry point
# --------------------------------------------------------------------------- #

def build_deck(results: dict, week_label: str, output_dir: Path) -> Path:
    """Build the full 16-slide deck and save to output_dir."""
    prs = _new_prs()

    channels = results.get("channels", {})
    videos = results.get("videos", [])
    engagement = results.get("engagement", {})

    _slide_01_cover(prs, week_label)
    _slide_02_key_findings(prs, results)
    _slide_03_channel_subscribers(prs, channels)
    _slide_04_channel_growth(prs, channels, results.get("prev_snapshot", {}))
    _slide_05_channel_table(prs, channels, engagement, videos)
    _slide_06_trending_topics(prs, results.get("trending_topics", Counter()))
    _slide_07_top_videos(prs, results.get("view_velocity", []))
    _slide_08_channel_engagement(prs, videos, channels, engagement)
    _slide_09_duration(prs, results.get("duration_engagement", {}))
    _slide_10_content_gaps(prs, results.get("gaps", []))
    _slide_11_recommendations(prs, results.get("gaps", []), results.get("trending_topics", Counter()))
    _slide_12_title_words(prs, results.get("title_words", Counter()))
    _slide_13_comment_sentiment(prs, results.get("sentiment", Counter()))
    _slide_14_thumbnail_styles(prs, videos)
    _slide_15_publish_timing(prs, results.get("timing", [[0]*24 for _ in range(7)]))
    _slide_16_upload_frequency(prs, videos, channels)

    output_dir.mkdir(parents=True, exist_ok=True)
    out_path = output_dir / f"{week_label}_niche_intelligence.pptx"
    prs.save(str(out_path))
    return out_path
